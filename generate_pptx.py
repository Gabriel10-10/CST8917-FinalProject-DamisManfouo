from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x00, 0x47, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
LIGHT = RGBColor(0xF2, 0xF6, 0xFF)


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_textbox(slide, title, 0.3, 0.15, 12.5, 1.0, size=28, bold=True, color=WHITE)


# ── Slide 1: Title ──────────────────────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, BLUE)
add_textbox(slide, "Expense Approval Workflow", 1, 1.5, 11, 1.2, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Azure Durable Functions  vs  Logic Apps + Service Bus", 1, 2.8, 11, 0.8, size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Damis Gabriel Manfouo  |  CST8917 — Serverless Applications  |  Algonquin College", 1, 4.2, 11, 0.6, size=16, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 2: Problem Statement ───────────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, LIGHT)
add_title_bar(slide, "Problem Statement")
bullets = [
    "• Automate employee expense approvals end-to-end",
    "• Support three outcomes: Auto-Approved, Manager Approved/Rejected, Escalated",
    "• Notify employees by email for every outcome",
    "• Compare two Azure serverless approaches for the same workflow",
]
for i, b in enumerate(bullets):
    add_textbox(slide, b, 0.5, 1.4 + i * 0.9, 12, 0.8, size=20)

# ── Slide 3: Version A — Durable Functions ──────────────────────────────────
slide = add_slide(prs)
set_bg(slide, LIGHT)
add_title_bar(slide, "Version A — Azure Durable Functions")
add_textbox(slide, "Architecture", 0.5, 1.3, 12, 0.5, size=18, bold=True, color=BLUE)
arch = "HTTP Client  →  Orchestrator  →  validate_expense  →  process_expense  →  notify_employee"
add_textbox(slide, arch, 0.5, 1.8, 12, 0.6, size=16)
add_textbox(slide, "Key Design Decisions", 0.5, 2.6, 12, 0.5, size=18, bold=True, color=BLUE)
points = [
    "• Human Interaction Pattern: race timer vs. ManagerDecision external event",
    "• Timer fires → status = Escalated (no manager response)",
    "• Manager calls POST /api/expense/{instanceId}/respond → approved or rejected",
    "• Real email via Azure Communication Services (ACS)",
]
for i, p in enumerate(points):
    add_textbox(slide, p, 0.5, 3.1 + i * 0.7, 12, 0.6, size=18)

# ── Slide 4: Version B — Logic Apps ─────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, LIGHT)
add_title_bar(slide, "Version B — Logic Apps + Service Bus")
add_textbox(slide, "Architecture", 0.5, 1.3, 12, 0.5, size=18, bold=True, color=BLUE)
arch = "Service Bus Queue  →  Logic App  →  Validation Function  →  Condition  →  expense-outcomes Topic"
add_textbox(slide, arch, 0.5, 1.8, 12, 0.6, size=16)
add_textbox(slide, "Key Design Decisions", 0.5, 2.6, 12, 0.5, size=18, bold=True, color=BLUE)
points = [
    "• Trigger: Service Bus queue 'expense-requests' (auto-complete)",
    "• Validation: Azure Function returns {valid, data} — Logic App branches on valid field",
    "• Manager approval: Send Approval Email (V2) with PT2M timeout",
    "• Scope wraps approve/reject branch; Send message 4 runs when Scope is skipped (timeout)",
    "• Outcomes published to 'expense-outcomes' topic with labels: approved / rejected / escalated",
]
for i, p in enumerate(points):
    add_textbox(slide, p, 0.5, 3.1 + i * 0.65, 12, 0.6, size=17)

# ── Slide 5: Comparison Table ────────────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, LIGHT)
add_title_bar(slide, "Comparison")

headers = ["Dimension", "Durable Functions", "Logic Apps"]
rows = [
    ["Dev Experience", "Code-first, full IDE support", "Visual designer, low-code"],
    ["Testability", "Local func start, unit tests", "Requires Azure, run history only"],
    ["Human Interaction", "External event + durable timer", "Send Approval Email + timeout"],
    ["Error Handling", "Full control via code", "Configure run after settings"],
    ["Observability", "App Insights, detailed logs", "Visual run history per step"],
    ["Cost (100/day)", "Low — consumption plan", "Low — per action pricing"],
]

col_widths = [2.8, 4.5, 4.5]
col_starts = [0.3, 3.2, 7.8]
row_height = 0.55
top_start = 1.3

for ci, (h, w) in enumerate(zip(headers, col_widths)):
    box = slide.shapes.add_shape(1, Inches(col_starts[ci]), Inches(top_start), Inches(w), Inches(row_height))
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.line.fill.background()
    add_textbox(slide, h, col_starts[ci] + 0.05, top_start + 0.05, w - 0.1, row_height - 0.1, size=15, bold=True, color=WHITE)

for ri, row in enumerate(rows):
    bg = RGBColor(0xE8, 0xF0, 0xFE) if ri % 2 == 0 else WHITE
    for ci, (cell, w) in enumerate(zip(row, col_widths)):
        box = slide.shapes.add_shape(1, Inches(col_starts[ci]), Inches(top_start + (ri + 1) * row_height), Inches(w), Inches(row_height))
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        add_textbox(slide, cell, col_starts[ci] + 0.05, top_start + (ri + 1) * row_height + 0.05, w - 0.1, row_height - 0.1, size=13)

# ── Slide 6: Challenges ──────────────────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, LIGHT)
add_title_bar(slide, "Key Challenges & Lessons Learned")
challenges = [
    ("Version A", [
        "• azure-durable-functions not found → correct package is azure-functions-durable",
        "• approval_task.result returns JSON string → must parse with json.loads()",
        "• Missing extensionBundle in host.json → durable bindings not registered",
    ]),
    ("Version B", [
        "• Logic Apps rejects custom routes → fix: @app.route(route='') with @app.function_name()",
        "• Service Bus delivers base64 content → Compose step needed to decode before function call",
        "• Condition boolean vs string: valid=True fails → must use equals(body,true) expression",
        "• Validation failure response had no data field → updated function to always return data",
    ]),
]
y = 1.4
for version, points in challenges:
    add_textbox(slide, version, 0.5, y, 12, 0.45, size=18, bold=True, color=BLUE)
    y += 0.45
    for p in points:
        add_textbox(slide, p, 0.5, y, 12, 0.5, size=16)
        y += 0.5
    y += 0.2

# ── Slide 7: Recommendation ──────────────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, LIGHT)
add_title_bar(slide, "Recommendation")
add_textbox(slide, "Choose Logic Apps when:", 0.5, 1.4, 12, 0.5, size=20, bold=True, color=BLUE)
for i, p in enumerate([
    "• Team prefers low-code / visual workflows",
    "• Rapid prototyping without writing infrastructure code",
    "• Built-in connectors (Outlook, Service Bus) cover the use case",
]):
    add_textbox(slide, p, 0.5, 1.9 + i * 0.6, 12, 0.5, size=18)

add_textbox(slide, "Choose Durable Functions when:", 0.5, 3.8, 12, 0.5, size=20, bold=True, color=BLUE)
for i, p in enumerate([
    "• Complex branching logic requiring full code control",
    "• Local testing and CI/CD pipelines are a priority",
    "• Team is comfortable with Python and Azure Functions SDK",
]):
    add_textbox(slide, p, 0.5, 4.3 + i * 0.6, 12, 0.5, size=18)

# ── Slide 8: Demo Video ──────────────────────────────────────────────────────
slide = add_slide(prs)
set_bg(slide, BLUE)
add_textbox(slide, "Demo Video", 1, 1.5, 11, 1, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "See video link in presentation/video-link.md", 1, 3, 11, 0.8, size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Damis Gabriel Manfouo  |  CST8917  |  Algonquin College", 1, 5.5, 11, 0.6, size=16, color=WHITE, align=PP_ALIGN.CENTER)

prs.save("presentation/CST8917-FinalProject-DamisManfouo.pptx")
print("Saved: presentation/CST8917-FinalProject-DamisManfouo.pptx")
